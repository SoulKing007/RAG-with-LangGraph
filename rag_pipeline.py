import streamlit as st
import os
import tempfile
import shutil
from pathlib import Path
from typing import List, Dict, Any, Optional
from typing_extensions import TypedDict, Annotated
import time
import requests
from urllib.parse import urlparse

# Core LangChain imports
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain import hub
from langgraph.graph import START, StateGraph
from langchain_core.prompts import PromptTemplate

# Gemini integration
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings

# ChromaDB integration
from langchain_chroma import Chroma

# Document loaders
from langchain_community.document_loaders import (
    PyPDFLoader, 
    Docx2txtLoader,
    WebBaseLoader,
    DirectoryLoader,
    TextLoader
)

# PowerPoint loader
from pptx import Presentation
import docx
import PyPDF2
import pdfplumber
import bs4

class RAGState(TypedDict):
    """State for RAG pipeline"""
    question: str
    context: List[Document]
    answer: str
    sources: List[str]

class DocumentProcessor:
    """Handles different document types"""
    
    @staticmethod
    def load_pdf(file_path: str) -> List[Document]:
        """Load PDF documents"""
        try:
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            return documents
        except Exception as e:
            st.error(f"Error loading PDF {file_path}: {str(e)}")
            return []
    
    @staticmethod
    def load_docx(file_path: str) -> List[Document]:
        """Load Word documents"""
        try:
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
            return documents
        except Exception as e:
            st.error(f"Error loading DOCX {file_path}: {str(e)}")
            return []
    
    @staticmethod
    def load_pptx(file_path: str) -> List[Document]:
        """Load PowerPoint presentations"""
        try:
            presentation = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                text_content.append(slide_text)
            
            full_text = "\n\n".join(text_content)
            doc = Document(
                page_content=full_text,
                metadata={"source": file_path, "type": "pptx"}
            )
            return [doc]
        except Exception as e:
            st.error(f"Error loading PPTX {file_path}: {str(e)}")
            return []
    
    @staticmethod
    def load_txt(file_path: str) -> List[Document]:
        """Load text files"""
        try:
            loader = TextLoader(file_path)
            documents = loader.load()
            return documents
        except Exception as e:
            st.error(f"Error loading TXT {file_path}: {str(e)}")
            return []
    
    @staticmethod
    def load_web_url(url: str) -> List[Document]:
        """Load content from web URL"""
        try:
            loader = WebBaseLoader(
                web_paths=[url],
                bs_kwargs=dict(
                    parse_only=bs4.SoupStrainer(
                        ["p", "h1", "h2", "h3", "h4", "h5", "h6", "article", "section"]
                    )
                )
            )
            documents = loader.load()
            return documents
        except Exception as e:
            st.error(f"Error loading URL {url}: {str(e)}")
            return []
    
    @staticmethod
    def load_directory(dir_path: str) -> List[Document]:
        """Load all supported files from directory"""
        documents = []
        supported_extensions = ['.pdf', '.docx', '.pptx', '.txt']
        
        for file_path in Path(dir_path).rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                file_str = str(file_path)
                if file_path.suffix.lower() == '.pdf':
                    documents.extend(DocumentProcessor.load_pdf(file_str))
                elif file_path.suffix.lower() == '.docx':
                    documents.extend(DocumentProcessor.load_docx(file_str))
                elif file_path.suffix.lower() == '.pptx':
                    documents.extend(DocumentProcessor.load_pptx(file_str))
                elif file_path.suffix.lower() == '.txt':
                    documents.extend(DocumentProcessor.load_txt(file_str))
        
        return documents

class RAGPipeline:
    """Complete RAG Pipeline with LangGraph"""
    
    def __init__(self, google_api_key: str, persist_directory: str = "./chroma_db"):
        """Initialize RAG pipeline"""
        self.google_api_key = google_api_key
        self.persist_directory = persist_directory
        
        # Initialize Gemini LLM
        self.llm = ChatGoogleGenerativeAI(
            model="gemini-2.5-flash",
            google_api_key=google_api_key,
            temperature=0.3
        )
        
        # Initialize embeddings
        self.embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001",
            google_api_key=google_api_key
        )
        
        # Initialize ChromaDB vector store
        self.vector_store = Chroma(
            persist_directory=persist_directory,
            embedding_function=self.embeddings
        )
        
        # Initialize text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            add_start_index=True
        )
        
        # Create RAG prompt
        self.rag_prompt = PromptTemplate.from_template("""
        You are an assistant for question-answering tasks. Use the following pieces of retrieved context to answer the question. 
        If you don't know the answer, just say that you don't know. Use three sentences maximum and keep the answer concise.
        
        Question: {question}
        Context: {context}
        Answer:
        """)
        
        # Build LangGraph
        self.graph = self._build_graph()
    
    def _build_graph(self):
        """Build LangGraph for RAG pipeline"""
        
        def retrieve(state: RAGState):
            """Retrieve relevant documents"""
            retrieved_docs = self.vector_store.similarity_search(
                state["question"], 
                k=4
            )
            sources = list(set([doc.metadata.get("source", "Unknown") for doc in retrieved_docs]))
            return {
                "context": retrieved_docs,
                "sources": sources
            }
        
        def generate(state: RAGState):
            """Generate answer using retrieved context"""
            docs_content = "\n\n".join(doc.page_content for doc in state["context"])
            messages = self.rag_prompt.invoke({
                "question": state["question"], 
                "context": docs_content
            })
            response = self.llm.invoke(messages)
            return {"answer": response.content}
        
        # Create graph
        graph_builder = StateGraph(RAGState).add_sequence([retrieve, generate])
        graph_builder.add_edge(START, "retrieve")
        return graph_builder.compile()
    
    def add_documents(self, documents: List[Document]) -> int:
        """Add documents to vector store"""
        if not documents:
            return 0
            
        # Split documents
        splits = self.text_splitter.split_documents(documents)
        
        # Add to vector store
        self.vector_store.add_documents(splits)
        
        return len(splits)
    
    def query(self, question: str) -> Dict[str, Any]:
        """Query the RAG pipeline"""
        result = self.graph.invoke({"question": question})
        return {
            "question": question,
            "answer": result["answer"],
            "sources": result.get("sources", []),
            "context": result.get("context", [])
        }
    
    def get_vector_store_stats(self) -> Dict[str, Any]:
        """Get statistics about the vector store"""
        try:
            collection = self.vector_store._collection
            count = collection.count()
            return {
                "document_count": count,
                "collection_name": collection.name if hasattr(collection, 'name') else "default"
            }
        except:
            return {"document_count": 0, "collection_name": "default"}

def main():
    """Streamlit UI"""
    st.set_page_config(
        page_title="RAG Pipeline",
        page_icon="🤖",
        layout="wide"
    )
    
    # Custom CSS for clean black and white UI
    st.markdown("""
    <style>.main {
    background-color: #ffffff;
    color: #000000;
}

.stSidebar {
    background-color: #ffffff;
    border-right: 1px solid #e0e0e0;
}

.stSidebar a {
    color: #000000;
}

.stSidebar .stButton > button {
    background-color: #ffffff;
    color: #000000;
    border: 1px solid #000000;
}

.stSidebar .stButton > button:hover {
    background-color: #f0f0f0;
    color: #000000;
}

.stButton > button {
    background-color: #ffffff;
    color: #000000;
    border: 1px solid #000000;
}

.stButton > button:hover {
    background-color: #f0f0f0;
    color: #000000;
}

.metric-container {
    background-color: #ffffff;
    padding: 0.75rem;
    border-radius: 0.25rem;
    border: 1px solid #e0e0e0;
    margin: 0.5rem 0;
}

.stTextInput > div > div > input {
    background-color: #ffffff;
    color: #000000;
    border: 1px solid #e0e0e0;
}

.stTextArea > div > div > textarea {
    background-color: #ffffff;
    color: #000000;
    border: 1px solid #e0e0e0;
}

.stExpander {
    border: 1px solid #e0e0e0;
    border-radius: 0.25rem;
}

.stMetric {
    background-color: #ffffff;
}
    </style>
    """, unsafe_allow_html=True)
    
    st.title("🤖 Complete RAG Pipeline")
    st.markdown("---")
    
    # Sidebar for configuration
    with st.sidebar:
        st.header("⚙️ Configuration")
        
        # Google API Key input
        google_api_key = st.text_input(
            "Google API Key",
            type="password",
            help="Enter your Google API key for Gemini"
        )
        
        if not google_api_key:
            st.warning("Please enter your Google API key to continue.")
            st.stop()
        
        st.markdown("---")
        
        # Initialize RAG pipeline
        if "rag_pipeline" not in st.session_state:
            with st.spinner("Initializing RAG pipeline..."):
                st.session_state.rag_pipeline = RAGPipeline(google_api_key)
        
        # Vector store statistics
        st.header("📊 Vector Store Stats")
        stats = st.session_state.rag_pipeline.get_vector_store_stats()
        
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Documents", stats["document_count"])
        with col2:
            st.metric("Collection", stats["collection_name"])
    
    # Main interface
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.header("📄 Document Upload")
        
        # File upload section
        uploaded_files = st.file_uploader(
            "Upload Documents",
            type=["pdf", "docx", "pptx", "txt"],
            accept_multiple_files=True,
            help="Upload PDF, Word, PowerPoint, or text files"
        )
        
        # Web URL section
        st.subheader("🌐 Web URL")
        web_url = st.text_input("Enter web URL to scrape:")
        
        # Directory upload simulation (using multiple files)
        st.subheader("📁 Multiple Files")
        st.info("Use the file uploader above to upload multiple files at once")
        
        # Process documents button
        if st.button("🔄 Process Documents", type="primary"):
            if uploaded_files or web_url:
                with st.spinner("Processing documents..."):
                    all_documents = []
                    
                    # Process uploaded files
                    if uploaded_files:
                        for uploaded_file in uploaded_files:
                            # Save uploaded file temporarily
                            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_file:
                                tmp_file.write(uploaded_file.getvalue())
                                tmp_path = tmp_file.name
                            
                            # Process based on file type
                            file_extension = uploaded_file.name.split('.')[-1].lower()
                            
                            if file_extension == 'pdf':
                                docs = DocumentProcessor.load_pdf(tmp_path)
                            elif file_extension == 'docx':
                                docs = DocumentProcessor.load_docx(tmp_path)
                            elif file_extension == 'pptx':
                                docs = DocumentProcessor.load_pptx(tmp_path)
                            elif file_extension == 'txt':
                                docs = DocumentProcessor.load_txt(tmp_path)
                            else:
                                docs = []
                            
                            # Update metadata
                            for doc in docs:
                                doc.metadata["source"] = uploaded_file.name
                                doc.metadata["type"] = file_extension
                            
                            all_documents.extend(docs)
                            
                            # Clean up temp file
                            os.unlink(tmp_path)
                    
                    # Process web URL
                    if web_url:
                        web_docs = DocumentProcessor.load_web_url(web_url)
                        all_documents.extend(web_docs)
                    
                    # Add documents to vector store
                    if all_documents:
                        num_chunks = st.session_state.rag_pipeline.add_documents(all_documents)
                        st.success(f"✅ Processed {len(all_documents)} documents into {num_chunks} chunks")
                        
                        # Show document details
                        st.subheader("📋 Processed Documents")
                        for i, doc in enumerate(all_documents[:5]):  # Show first 5
                            with st.expander(f"Document {i+1}: {doc.metadata.get('source', 'Unknown')}"):
                                st.write(f"**Type:** {doc.metadata.get('type', 'Unknown')}")
                                st.write(f"**Content Preview:**")
                                st.write(doc.page_content[:500] + "..." if len(doc.page_content) > 500 else doc.page_content)
                        
                        if len(all_documents) > 5:
                            st.info(f"... and {len(all_documents) - 5} more documents")
                        
                        # Update stats
                        st.rerun()
                    else:
                        st.error("❌ No documents were processed successfully")
            else:
                st.warning("Please upload files or enter a web URL")
    
    with col2:
        st.header("💬 Query Interface")
        
        # Query input
        user_question = st.text_area(
            "Ask a question about your documents:",
            height=100,
            placeholder="What is the main topic discussed in the documents?"
        )
        
        # Query button
        if st.button("🔍 Search", type="primary"):
            if user_question and st.session_state.rag_pipeline.get_vector_store_stats()["document_count"] > 0:
                with st.spinner("Searching for answer..."):
                    result = st.session_state.rag_pipeline.query(user_question)
                    
                    # Display results
                    st.subheader("🎯 Answer")
                    st.write(result["answer"])
                    
                    # Display sources
                    if result["sources"]:
                        st.subheader("📚 Sources")
                        for source in result["sources"]:
                            st.write(f"• {source}")
                    
                    # Display retrieved context
                    if result["context"]:
                        st.subheader("📄 Retrieved Context")
                        for i, doc in enumerate(result["context"]):
                            with st.expander(f"Context {i+1} from {doc.metadata.get('source', 'Unknown')}"):
                                st.write(doc.page_content)
            
            elif not user_question:
                st.warning("Please enter a question")
            else:
                st.warning("Please upload and process documents first")
        
        # Query history
        if "query_history" not in st.session_state:
            st.session_state.query_history = []
        
        if st.session_state.query_history:
            st.subheader("📜 Recent Queries")
            for i, (question, answer) in enumerate(st.session_state.query_history[-3:]):
                with st.expander(f"Q: {question[:50]}..."):
                    st.write(f"**A:** {answer}")
    
    # Footer
    st.markdown("---")
    st.markdown(
        """
        <div style='text-align: center; color: #666666;'>
            <p>🤖 RAG Pipeline powered by LangGraph, Gemini & ChromaDB</p>
        </div>
        """,
        unsafe_allow_html=True
    )

if __name__ == "__main__":

    main()
